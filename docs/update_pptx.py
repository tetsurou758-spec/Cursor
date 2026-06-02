# -*- coding: utf-8 -*-
"""
AI協働開発レポート_v3_白基調.pptx を更新するスクリプト
5/29以降の実装成果（契約照会・コンタクト履歴・満期管理・保険金支払状況）を反映
"""

import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

SRC = r"C:\Users\yoshi\OneDrive\ドキュメント\Cursor\docs\AI協働開発レポート_v3_白基調.pptx"
DST_V3 = SRC
DST_V4 = r"C:\Users\yoshi\OneDrive\ドキュメント\Cursor\docs\AI協働開発レポート_v4_白基調.pptx"

# ============================================================
# ヘルパー関数
# ============================================================

def rgb(hex_str):
    """'RRGGBB' → RGBColor"""
    return RGBColor.from_string(hex_str)


def add_solid_rect(slide, left, top, width, height, fill_hex, line_hex=None, radius=None):
    """塗りつぶし矩形を追加。rounded rectangle も可（radius は EMU 単位）。"""
    from pptx.util import Emu
    from pptx.oxml import parse_xml
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = MSO_SHAPE_TYPE.AUTO_SHAPE
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    if line_hex:
        shape.line.color.rgb = rgb(line_hex)
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_name="Meiryo UI", font_size=None, bold=False,
                color_hex="0F172A", align=PP_ALIGN.LEFT,
                wrap=True, line_spacing=None):
    """テキストボックスを追加し整形。"""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    if font_name:
        font.name = font_name
    if font_size:
        font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = rgb(color_hex)
    return txb


def add_multiline_textbox(slide, lines, left, top, width, height,
                          font_name="Meiryo UI", font_size=10, bold=False,
                          color_hex="475569", align=PP_ALIGN.LEFT):
    """複数行テキストボックスを追加（行ごとにパラグラフ）。"""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        font = run.font
        font.name = font_name
        font.size = Pt(font_size)
        font.bold = bold
        font.color.rgb = rgb(color_hex)
    return txb


def clone_slide(prs, slide_idx):
    """
    指定インデックスのスライドを複製して末尾に追加し、新スライドを返す。
    """
    template = prs.slides[slide_idx]
    # スライドレイアウトは同じものを使う
    slide_layout = template.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)
    # テンプレートの XML ツリーをコピー
    sp_tree = template.shapes._spTree
    new_sp_tree = new_slide.shapes._spTree
    # 既存シェイプを削除
    for sp in new_sp_tree.findall('.//' + qn('p:sp')):
        sp.getparent().remove(sp)
    for sp in new_sp_tree.findall('.//' + qn('p:pic')):
        sp.getparent().remove(sp)
    for sp in new_sp_tree.findall('.//' + qn('p:graphicFrame')):
        sp.getparent().remove(sp)
    # テンプレートの全シェイプを新スライドにコピー
    for elem in sp_tree:
        tag = elem.tag
        if tag in (qn('p:sp'), qn('p:pic'), qn('p:graphicFrame'), qn('p:grpSp'), qn('p:cxnSp')):
            new_sp_tree.append(copy.deepcopy(elem))
    return new_slide


# ============================================================
# スライド別の更新処理
# ============================================================

def update_slide1_title(slide):
    """スライド1：タイトル数字（1.5日→継続中、8画面→12画面）を更新。"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text == '8':
                    run.text = '12'
                    return


def update_slide5_system(slide):
    """
    スライド5：システム構成スライド
    - FRONTENDに contract.html / contact.html / claim.html を追加
    - API SERVERに /api/contracts, /api/contacts, /api/accidents を追加
    - DATABASEに contacts, contract_details, accident_payments を追加
    既存テキストシェイプを直接書き換える。
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()

        # admin.html のボックスを「contract.html / 契約照会」に更新
        if 'admin.html' in text and '権限管理' in text:
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    if 'admin.html' in run.text:
                        run.text = 'contract.html'
                    if '権限管理' in run.text:
                        run.text = '契約照会'

        # /api/users のボックスを /api/contracts に更新
        if '/api/users' in text:
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    if '/api/users' in run.text:
                        run.text = '/api/contracts'
                    if 'ユーザー管理' in run.text:
                        run.text = '契約・詳細'

        # login_history のボックスを contacts に更新
        if 'login_history' in text and 'ログ' in text:
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    if 'login_history' in run.text:
                        run.text = 'contacts'
                    if 'ログ' in run.text:
                        run.text = 'コンタクト'


def update_slide6_screens(slide):
    """
    スライド6：完成画面一覧
    - 顧客管理の 'WIP' バッジを 'SHARED' に変更
    - 下段3カードに contract / contact / claim を追加（既存カードを更新）
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # WIP → SHARED へ変更
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text == 'WIP':
                    run.text = 'SHARED'

    # WIPバッジ背景色（D97706）を 4F46E5 に変更
    for shape in slide.shapes:
        if not shape.has_text_frame:
            try:
                if shape.fill.type and shape.fill.type.name == 'SOLID':
                    if str(shape.fill.fore_color.rgb) == 'D97706':
                        # WIP バッジの背景（小さいボックス）
                        if shape.width < 700000 and shape.top > 2800000:
                            shape.fill.fore_color.rgb = rgb('4F46E5')
                            shape.line.color.rgb = rgb('4F46E5')
            except Exception:
                pass


def update_slide8_metrics(slide):
    """
    スライド8：定量実績
    - 完成画面数 8 → 12
    - DBテーブル数 9 → 14
    - テストフォルダ 5 → 7
    - 自動テストケース 13 → 25+
    """
    replacements = {
        '8': '12',
        '9': '14',
        '5': '7',
        '13': '25+',
        '代理店・社員テーマ統合対応': '契約照会・コンタクト・保険金含む',
        '1,465件のダミーデータ投入': 'contacts等新規5TBL追加',
        '仕様書+証跡+Playwright自動化': '仕様書+証跡+Playwright自動化',
        'XSS・SQLiセキュリティ含む': 'XSS・SQLi・E2Eシナリオ含む',
    }
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text in replacements:
                    run.text = replacements[run.text]


# ============================================================
# 新規スライドの作成
# ============================================================

def make_slide_header(slide, num_str, title_text, accent_hex='7C3AED'):
    """スライド共通ヘッダー（番号バッジ＋タイトル）を追加する。"""
    # 番号バッジ（背景）
    badge_bg = add_solid_rect(slide, 320040, 274320, 292608, 292608, accent_hex, accent_hex)
    # 番号テキスト
    add_textbox(slide, num_str, 320040, 274320, 292608, 292608,
                font_size=9, bold=True, color_hex='FFFFFF', align=PP_ALIGN.CENTER)
    # 区切り線
    add_solid_rect(slide, 320040, 658368, 8503920, 27432, 'E2E8F0')
    # タイトル
    add_textbox(slide, title_text, 749808, 274320, 8046720, 292608,
                font_size=20, bold=True, color_hex='0F172A')


def add_new_slide_contract_contact(prs):
    """
    新規スライド「06-2 / 06-3：契約照会・コンタクト履歴・保険金支払状況」
    既存スライド6のカードレイアウトを参考に3カード構成で作成。
    """
    blank_layout = prs.slide_layouts[6]  # blank レイアウト
    slide = prs.slides.add_slide(blank_layout)

    # 背景を白にする（デフォルトで白だが明示）
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb('FFFFFF')

    # ---- ヘッダー ----
    # 番号バッジ背景
    add_solid_rect(slide, 320040, 274320, 292608, 292608, '0891B2', '0891B2')
    add_textbox(slide, '06+', 320040, 274320, 292608, 292608,
                font_size=9, bold=True, color_hex='FFFFFF', align=PP_ALIGN.CENTER)
    # 区切り線
    add_solid_rect(slide, 320040, 658368, 8503920, 27432, 'E2E8F0')
    # タイトル
    add_textbox(slide, '新規実装画面 ── 契約照会 / コンタクト履歴 / 保険金支払状況',
                749808, 274320, 8046720, 292608,
                font_size=18, bold=True, color_hex='0F172A')

    # カード共通サイズ
    CARD_W = 2743200
    CARD_H = 3584160  # スライド高さいっぱいに近く
    TOP_BAR_H = 45720
    BADGE_W = 594360
    BADGE_H = 219456
    TITLE_H = 347472
    DIVIDER_H = 18288
    BODY_H = 1874520
    ROW_TOP = 750000  # カード開始位置（Y）

    cards = [
        {
            'left': 256032,
            'accent': '0891B2',
            'bg': 'ECFEFF',
            'badge_text': 'SHARED',
            'badge_bg': '4F46E5',
            'title': '契約照会画面',
            'title_color': '0891B2',
            'lines': [
                'contract.html / contract_detail.html',
                '証券番号・顧客名・種目で検索',
                '保険種目別詳細情報テーブル表示',
                '1件ヒット→詳細、複数→一覧遷移',
                'FastAPI: contract_router.py',
            ],
        },
        {
            'left': 3182112,
            'accent': '059669',
            'bg': 'ECFDF5',
            'badge_text': 'SHARED',
            'badge_bg': '4F46E5',
            'title': 'コンタクト履歴画面',
            'title_color': '059669',
            'lines': [
                'contact.html（2026-06-01実装）',
                'contactsテーブル新設',
                '履歴参照・編集（CRUD完備）',
                '満期管理「落ち」選択で自動登録',
                'ダッシュボードウィジェット動的化',
            ],
        },
        {
            'left': 6108192,
            'accent': 'E11D48',
            'bg': 'FFF1F2',
            'badge_text': 'SHARED',
            'badge_bg': '4F46E5',
            'title': '保険金支払状況画面',
            'title_color': 'E11D48',
            'lines': [
                'claim.html / claim_detail.html',
                'accidentsテーブル拡張',
                'accident_paymentsテーブル新設',
                '支払履歴・状況管理',
                'API: /api/accidents 系',
            ],
        },
    ]

    for card in cards:
        L = card['left']
        T = ROW_TOP

        # カード背景
        add_solid_rect(slide, L, T, CARD_W, 3400000, card['bg'], card['accent'])
        # 上部アクセントバー
        add_solid_rect(slide, L, T, CARD_W, TOP_BAR_H, card['accent'], card['accent'])
        # バッジ背景
        badge_left = L + CARD_W - BADGE_W - 45720
        add_solid_rect(slide, badge_left, T + 91440, BADGE_W, BADGE_H, card['badge_bg'], card['badge_bg'])
        add_textbox(slide, card['badge_text'],
                    badge_left, T + 91440, BADGE_W, BADGE_H,
                    font_size=8, bold=True, color_hex='FFFFFF', align=PP_ALIGN.CENTER)
        # タイトル
        add_textbox(slide, card['title'],
                    L + 109728, T + 91440, 1920240, TITLE_H,
                    font_size=13, bold=True, color_hex=card['title_color'])
        # 区切り線
        add_solid_rect(slide, L + 109728, T + 493920, 2523744, DIVIDER_H, 'E2E8F0')
        # 本文
        add_multiline_textbox(slide, card['lines'],
                              L + 109728, T + 566640, 2523744, 2200000,
                              font_size=10, color_hex='475569')

    return slide


def add_new_slide_timeline_extended(prs, existing_slide_idx=3):
    """
    スライド4（タイムライン）のDay2夕以降を拡張した新スライドを追加する。
    """
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb('FFFFFF')

    # ---- ヘッダー ----
    add_solid_rect(slide, 320040, 274320, 292608, 292608, '4F46E5', '4F46E5')
    add_textbox(slide, '04+', 320040, 274320, 292608, 292608,
                font_size=9, bold=True, color_hex='FFFFFF', align=PP_ALIGN.CENTER)
    add_solid_rect(slide, 320040, 658368, 8503920, 27432, 'E2E8F0')
    add_textbox(slide, '開発タイムライン続き ── 5/29以降の実装成果',
                749808, 274320, 8046720, 292608,
                font_size=20, bold=True, color_hex='0F172A')

    # タイムラインアイテム
    items = [
        ('Day 3  前半', '契約照会機能', '0891B2',
         'contract.html / contract_detail.html\n証券番号・種目別詳細テーブル・ダイレクト振り分けロジック'),
        ('Day 3  後半', '満期管理インライン編集', 'D97706',
         'ポップオーバー編集（メモ・フォローコール・更改STS）\n「落ち」選択でコンタクト履歴自動登録'),
        ('Day 4  前半', 'コンタクト履歴機能', '059669',
         'contactsテーブル新設・CRUD API完備\nダッシュボードウィジェット動的化'),
        ('Day 4  後半', '保険金支払状況機能', 'E11D48',
         'claim.html / claim_detail.html\naccident_paymentsテーブル・支払履歴管理'),
    ]

    # タイムライン縦線
    add_solid_rect(slide, 1280160, 750000, 36576, 3700000, 'E2E8F0')

    for i, (day_label, title, color, desc) in enumerate(items):
        item_top = 800000 + i * 950000
        # 丸ドット
        dot = slide.shapes.add_shape(9, 1243584, item_top, 109728, 109728)  # 9 = OVAL
        dot.fill.solid()
        dot.fill.fore_color.rgb = rgb(color)
        dot.line.fill.background()
        # 日付ラベル
        add_textbox(slide, day_label, 256032, item_top - 20000, 950000, 164592,
                    font_size=9, bold=True, color_hex='64748B')
        # タイトル
        add_textbox(slide, title, 1426464, item_top, 3200000, 200000,
                    font_size=14, bold=True, color_hex=color)
        # 説明
        add_textbox(slide, desc, 1426464, item_top + 200000, 7000000, 700000,
                    font_size=10, color_hex='475569')

    return slide


# ============================================================
# メイン処理
# ============================================================

def main():
    prs = Presentation(SRC)

    print("既存スライドを更新中...")

    # スライド1：タイトル（完成画面数更新）
    update_slide1_title(prs.slides[0])
    print("  スライド1 更新完了（完成画面数 8→12）")

    # スライド5：システム構成
    update_slide5_system(prs.slides[4])
    print("  スライド5 更新完了（FRONTEND/API/DB テキスト更新）")

    # スライド6：完成画面一覧（WIP→SHARED）
    update_slide6_screens(prs.slides[5])
    print("  スライド6 更新完了（顧客管理 WIP→SHARED）")

    # スライド8：定量実績
    update_slide8_metrics(prs.slides[7])
    print("  スライド8 更新完了（各指標更新）")

    print("新規スライドを追加中...")

    # 新規スライド：タイムライン続き
    add_new_slide_timeline_extended(prs, existing_slide_idx=3)
    print("  タイムライン拡張スライド追加完了")

    # 新規スライド：新規実装画面
    add_new_slide_contract_contact(prs)
    print("  新規実装画面スライド追加完了")

    # 保存
    prs.save(DST_V3)
    print(f"\n上書き保存完了: {DST_V3}")

    prs.save(DST_V4)
    print(f"v4 別名保存完了: {DST_V4}")

    print(f"\n総スライド数: {len(prs.slides)}")


if __name__ == '__main__':
    main()
